"""
Jornada de Compra — Preparação de Dados.

Upload de arquivos de eye-tracking + formulário de contexto do projeto.
"""

import streamlit as st
from utils import auth, ui
from utils.icons import page_title

auth.require_module("jornada_compra")

from io import BytesIO
from pptx import Presentation

from utils.jornada_loader import (
    load_jornada_from_upload,
    get_jornada_summary,
    get_jornada_participants,
    get_jornada_marcas,
)
from utils.organization_data import hydrate_session_state, save_session_state


_JORNADA_STATE_KEYS = (
    "jc_data",
    "jc_pptx_text",
    "jc_projeto",
    "jc_entrevistas_summary",
)


def _extract_pptx_text(file_bytes: bytes) -> str:
    """Extrai todo o texto de um arquivo .pptx."""
    prs = Presentation(BytesIO(file_bytes))
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells
                    )
                    if row_text.replace(" | ", "").strip():
                        parts.append(row_text)
        if parts:
            slides_text.append(f"--- Slide {i} ---\n" + "\n".join(parts))
    return "\n\n".join(slides_text)


hydrate_session_state("jornada_compra", _JORNADA_STATE_KEYS)

ui.inject_theme()
ui.breadcrumb("Jornada de Compra", "Preparação de Dados")
page_title(
    "folder-open",
    "Preparação de Dados",
    "Eye tracking, entrevistas e contexto do projeto.",
)

# ==================================================================
# Seção 1: Upload de Dados
# ==================================================================
st.subheader("Upload de Dados")
st.markdown(
    "Envie os arquivos de eye-tracking do estudo. "
    "**Todos os arquivos são opcionais** — carregue os que estiverem disponíveis."
)

col1, col2, col3 = st.columns(3)

with col1:
    tabelas_file = st.file_uploader(
        "Tabelas (per-participante × AOI)",
        type=["csv", "xlsx"],
        key="jc_tabelas",
        help="Banco_Tabelas.csv — dados brutos por participante e AOI",
    )
    por_marca_file = st.file_uploader(
        "Por Marca (agrupado por marca)",
        type=["csv", "xlsx"],
        key="jc_por_marca",
        help="Banco_PorMarca.csv — dados agrupados por marca com metadados",
    )

with col2:
    medias_file = st.file_uploader(
        "Médias (médias por AOI)",
        type=["csv", "xlsx"],
        key="jc_medias",
        help="Banco_medias.csv — médias das métricas por AOI",
    )
    visual_share_file = st.file_uploader(
        "Visual Share (share por marca)",
        type=["csv", "xlsx"],
        key="jc_visual_share",
        help="Banco_TBVisualShare.csv — share visual por marca",
    )

with col3:
    anova_file = st.file_uploader(
        "ANOVA (testes estatísticos)",
        type=["csv", "xlsx"],
        key="jc_anova",
        help="Banco_ANOVA.csv — ANOVA por métrica",
    )
    consolidado_file = st.file_uploader(
        "Consolidado (.xlsx)",
        type=["xlsx"],
        key="jc_consolidado",
        help="Banco_Consolidado.xlsx — Excel consolidado",
    )

entrevistas_file = st.file_uploader(
    "Entrevistas (transcrições qualitativas)",
    type=["csv", "xlsx"],
    key="jc_entrevistas",
    help="Entrevistas.csv — colunas: arquivo, ep, identificacao, texto",
)

relatorio_file = st.file_uploader(
    "Relatório PPTX (contexto adicional para análise de IA)",
    type=["pptx"],
    key="jc_relatorio_pptx",
    help="Relatório em PowerPoint — o texto será extraído e usado como contexto para a IA",
)

if relatorio_file is not None:
    pptx_text = _extract_pptx_text(relatorio_file.getvalue())
    st.session_state["jc_pptx_text"] = pptx_text
    save_session_state("jornada_compra", _JORNADA_STATE_KEYS)

# Carregar se algum arquivo foi enviado
any_file = any([
    tabelas_file, por_marca_file, medias_file,
    visual_share_file, anova_file, consolidado_file,
    entrevistas_file,
])

if any_file:
    data = load_jornada_from_upload(
        tabelas_file=tabelas_file,
        por_marca_file=por_marca_file,
        medias_file=medias_file,
        visual_share_file=visual_share_file,
        anova_file=anova_file,
        consolidado_file=consolidado_file,
        entrevistas_file=entrevistas_file,
    )
    st.session_state["jc_data"] = data
    save_session_state("jornada_compra", _JORNADA_STATE_KEYS)

    # Avisos
    if "_errors" in data:
        for err in data["_errors"]:
            st.warning(err)

# ------------------------------------------------------------------
# Resumo dos dados carregados
# ------------------------------------------------------------------
data = st.session_state.get("jc_data", {})
loaded_keys = [k for k in data if k != "_errors"]

if loaded_keys:
    st.divider()
    st.subheader("Resumo dos Dados")

    summary = get_jornada_summary(data)

    n_entrevistas = len(data["entrevistas"]) if "entrevistas" in data else 0
    cols = st.columns(5)
    cols[0].metric("Arquivos carregados", summary["n_arquivos"])
    cols[1].metric("Participantes", summary["n_participantes"])
    cols[2].metric("AOIs", summary["n_aois"])
    cols[3].metric("Marcas", summary["n_marcas"])
    cols[4].metric("Entrevistas", n_entrevistas)

    # Listar arquivos carregados
    st.markdown("**Arquivos disponíveis:** " + ", ".join(
        f"`{k}`" for k in loaded_keys
    ))

    # Participantes e marcas
    participants = get_jornada_participants(data)
    marcas = get_jornada_marcas(data)

    if participants:
        with st.expander(f"Participantes ({len(participants)})"):
            st.write(", ".join(participants))

    if marcas:
        with st.expander(f"Marcas ({len(marcas)})"):
            st.write(", ".join(marcas))

    pptx_text = st.session_state.get("jc_pptx_text", "")
    if pptx_text:
        with st.expander("Conteúdo extraído do PPTX"):
            st.text(pptx_text[:3000] + ("\n..." if len(pptx_text) > 3000 else ""))

    # Preview de cada arquivo
    st.divider()
    st.subheader("Pré-visualização")

    labels = {
        "tabelas": "Tabelas (per-participante)",
        "consolidado": "Consolidado",
        "por_marca": "Por Marca",
        "medias": "Médias",
        "visual_share": "Visual Share",
        "anova": "ANOVA",
        "entrevistas": "Entrevistas",
    }
    for key in loaded_keys:
        df = data[key]
        label = labels.get(key, key)
        if key == "entrevistas" and "texto" in df.columns:
            with st.expander(f"{label} — {len(df)} entrevistas"):
                id_col = next(
                    (c for c in df.columns if c.lower() in ("identificacao", "identificação")),
                    df.columns[0],
                )
                for _, row in df.iterrows():
                    st.markdown(f"**{row.get(id_col, '')}**")
                    st.text(str(row["texto"])[:500] + ("…" if len(str(row["texto"])) > 500 else ""))
                    st.divider()
        else:
            with st.expander(f"{label} — {len(df)} linhas, {len(df.columns)} colunas"):
                st.dataframe(df.head(50), width='stretch')

else:
    st.info("Carregue pelo menos um arquivo para começar.")

# ==================================================================
# Seção 2: Contexto do Projeto
# ==================================================================
st.divider()
st.subheader("Contexto do Projeto")
st.markdown(
    "Preencha as informações abaixo para contextualizar a análise de IA. "
    "Essas informações serão usadas na página de **Análise** para gerar "
    "insights mais relevantes."
)

# Carregar valores salvos
projeto = st.session_state.get("jc_projeto", {})

with st.form("projeto_form", clear_on_submit=False):
    nome_projeto = st.text_input(
        "Nome do Projeto",
        value=projeto.get("nome", ""),
        placeholder="",
    )

    especialidade = st.text_area(
        "1. Como você quer que eu pense? Diga um pouco em qual a área de "
        "especialidade que vamos discutir.",
        value=projeto.get("especialidade", ""),
        placeholder=(
            "Ex: Somos uma equipe de neuromarketing focada em comportamento "
            "do consumidor no ponto de venda (PDV) para produtos de higiene "
            "feminina..."
        ),
        height=120,
    )

    historico = st.text_area(
        "2. Me conta um pouco sobre o histórico do problema que vamos discutir.",
        value=projeto.get("historico", ""),
        placeholder=(
            "Ex: A marca Intimus está reavaliando sua estratégia de embalagens "
            "e posicionamento na gôndola após mudanças de portfólio..."
        ),
        height=120,
    )

    problemas = st.text_area(
        "3. Quais os problemas centrais que o estudo deve responder?",
        value=projeto.get("problemas", ""),
        placeholder=(
            "Ex: 1) Qual marca atrai mais atenção visual? "
            "2) As novas embalagens da Intimus são mais notadas? "
            "3) Quais produtos são mais fixados vs. apenas olhados?"
        ),
        height=120,
    )

    submitted = st.form_submit_button("Salvar Contexto", type="primary")

    if submitted:
        st.session_state["jc_projeto"] = {
            "nome": nome_projeto,
            "especialidade": especialidade,
            "historico": historico,
            "problemas": problemas,
        }
        save_session_state("jornada_compra", _JORNADA_STATE_KEYS)
        st.success("Contexto do projeto salvo com sucesso!")

# ==================================================================
# Navegação por etapas
# ==================================================================
st.divider()
col_nav1, col_nav2 = st.columns(2)

with col_nav1:
    if st.button("Voltar para o Início", width='stretch'):
        st.session_state.pop("modulo", None)
        st.switch_page("home.py")

with col_nav2:
    if st.button("Avançar para Análise", width='stretch', type="primary"):
        st.switch_page("modules/jornada_compra/analise.py")
